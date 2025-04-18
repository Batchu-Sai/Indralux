
import io
import os
from pptx import Presentation
from PIL import Image
import io

def extract_clean_images_from_pptx(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    extracted = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not shape.shape_type == 13:  # Skip if not a picture
                continue
            if not hasattr(shape, "image"):
                continue

            image = shape.image
            ext = image.ext
            name = f"slide{slide_idx:02d}_img{shape_idx:02d}.{ext}"
            path = os.path.join(output_dir, name)

            # Extract all images, skip none
            try:
                img_data = io.BytesIO(image.blob)
                with Image.open(img_data) as im:
                    im.save(path)
                    extracted.append(name)
            except Exception as e:
                print(f"Skipping image: {e}")
                continue

    return extracted
