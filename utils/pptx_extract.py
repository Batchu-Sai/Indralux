from pptx import Presentation
import os
from PIL import Image
from io import BytesIO

def extract_clean_images_from_pptx(pptx_path, output_dir, min_width=300, min_height=300):
    """
    Extracts high-quality images from slides, ignoring small/placeholder graphics.

    Parameters:
        pptx_path (str): PowerPoint file path
        output_dir (str): Destination folder
        min_width (int): Minimum image width to include
        min_height (int): Minimum image height to include

    Returns:
        List of saved image filenames
    """
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    extracted = []

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.shape_type == 13 or not hasattr(shape, "image"):
                continue

            image = shape.image
            try:
                img = Image.open(BytesIO(image.blob)).convert("RGB")
            except Exception:
                continue  # Skip corrupt or unreadable images

            if img.width < min_width or img.height < min_height:
                continue  # Skip tiny artifacts or logos

            name = f"slide{i:02d}_clean.png"
            path = os.path.join(output_dir, name)
            img.save(path)
            extracted.append(name)

    return extracted
